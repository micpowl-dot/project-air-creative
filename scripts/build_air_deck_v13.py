"""Project AIR PowerPoint template v4.

v3 + real Figma slot exports (squiggle, quarter-circles, piano-stripes, dots,
top-edge curve). Replaces the bad circle-based approximations.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
import os

# --- Tokens (poster palette from Figma 324:1019) ---
ANCHOR_DARK = RGBColor(0x29, 0x29, 0x29)
ANCHOR_LIGHT = RGBColor(0xF0, 0xF0, 0xEB)
DARK_BOX_FILL = RGBColor(0x46, 0x12, 0x5B)  # deep violet from poster

# Note: "cyan" key now holds royal blue (#0062ff) — key kept for code stability
PALETTES = {
    "magenta": {"dominant": RGBColor(0xFB, 0x00, 0xFF), "accent": RGBColor(0xF4, 0xDC, 0x52)},
    "cyan":    {"dominant": RGBColor(0x00, 0x62, 0xFF), "accent": RGBColor(0x67, 0xFA, 0xE0)},
    "violet":  {"dominant": RGBColor(0x46, 0x12, 0x5B), "accent": RGBColor(0xF4, 0xDC, 0x52)},
}

DISPLAY_FONT = "Fira Sans Extra Condensed"
MONO_FONT = "Fira Code"

PNG_DIR = "/Users/lee.payne@weather.com/Dev/project-air-creative/assets/templates/slot-elements/png"
SLOT_SQUIGGLE = f"{PNG_DIR}/squiggle.png"
SLOT_TOP_CURVE = f"{PNG_DIR}/top-edge-curve.png"
SLOT_QUARTER_CIRCLES = f"{PNG_DIR}/quarter-circles.png"
SLOT_STRIPES_LIGHT = f"{PNG_DIR}/Vector1_light.png"
SLOT_STRIPES_DARK = f"{PNG_DIR}/Vector1_dark.png"
SLOT_DOTS_LIGHT = f"{PNG_DIR}/Pattern3_light.png"
SLOT_DOTS_DARK = f"{PNG_DIR}/Pattern3_dark.png"
SLOT_CONVERGING_LIGHT = f"{PNG_DIR}/converging_light.png"
SLOT_CONVERGING_DARK = f"{PNG_DIR}/converging_dark.png"
SLOT_RING_PATTERN_LIGHT = f"{PNG_DIR}/Pattern_light.png"
SLOT_RING_PATTERN_DARK = f"{PNG_DIR}/Pattern_dark.png"
# Real Figma logos — one per palette accent
LOGO_AIDAY = {k: f"{PNG_DIR}/ai-day-mark_{k}.png"
              for k in ["magenta", "cyan", "violet"]}
LOGO_AIR = {k: f"{PNG_DIR}/air-lockup_{k}.png"
            for k in ["magenta", "cyan", "violet"]}
QR_RSVP = f"{PNG_DIR}/qr.png"

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text,
             font=DISPLAY_FONT, size=14, color=ANCHOR_DARK, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    return tb


def add_text_runs(slide, left, top, width, height, runs,
                  font=DISPLAY_FONT, size=14, color=ANCHOR_DARK, bold=False,
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    for text, override in runs:
        r = p.add_run()
        r.text = text
        f = r.font
        f.name = font
        f.size = Pt(size)
        f.bold = bold
        f.color.rgb = override if override is not None else color
    return tb


def no_shadow(shape):
    """Suppress theme-default drop shadow on a shape by inserting empty effectLst."""
    spPr = shape._element.spPr
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    etree.SubElement(spPr, qn('a:effectLst'))


def add_rect(slide, left, top, width, height, color, line_color=None, shadow=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    no_shadow(shp)
    return shp


def add_oval(slide, left, top, width, height, fill_color, line_color=None, line_w=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    if fill_color is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        if line_w is not None:
            shp.line.width = line_w
    no_shadow(shp)
    return shp


def add_ring_badge(slide, cx, cy, rings, max_diam, stroke_w_pt=4):
    n = len(rings)
    for i, color in enumerate(rings):
        scale = (n - i) / n
        d = int(max_diam * scale)
        add_oval(slide, cx - d // 2, cy - d // 2, d, d,
                 fill_color=None, line_color=color, line_w=Pt(stroke_w_pt))


def add_air_lockup(slide, left, top, height, accent_color, anchor_color=ANCHOR_LIGHT):
    word_h = int(height * 0.7)
    add_text_runs(
        slide, left, top, Inches(3), word_h,
        runs=[("a", anchor_color), ("i", accent_color), ("r", anchor_color)],
        font=DISPLAY_FONT, size=88, bold=True,
    )
    tag_y = top + word_h - Inches(0.05)
    add_text_runs(
        slide, left, tag_y, Inches(2.5), Inches(0.25),
        runs=[("AI", anchor_color), ("IN", accent_color), ("REACH", anchor_color)],
        font=MONO_FONT, size=9, bold=True,
    )


def add_date_with_accent(slide, left, top, accent_color,
                         base_color=ANCHOR_DARK, size=16):
    add_text_runs(
        slide, left, top, Inches(4), Inches(0.4),
        runs=[("JUNE.9.20", base_color), ("26", accent_color)],
        font=MONO_FONT, size=size, bold=True,
    )


def add_eyebrow(slide, left, top, text, color=ANCHOR_DARK):
    add_text(slide, left, top, Inches(4), Inches(0.2),
             text, font=MONO_FONT, size=9, color=color, bold=False)


def add_image(slide, path, left, top, width=None, height=None):
    return slide.shapes.add_picture(path, left, top, width=width, height=height)


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ============================================================
# SLIDE 1 — Title (Magenta) — with real squiggle + piano-stripes
# ============================================================
p = PALETTES["magenta"]
s = prs.slides.add_slide(BLANK)
set_bg(s, p["dominant"])
# Real piano-stripes (anchor-dark recolor) top-right
add_image(s, SLOT_STRIPES_DARK, Inches(6.5), Inches(0.0),
          width=Inches(3.5), height=Inches(2))
# Date
add_date_with_accent(s, Inches(0.4), Inches(0.4),
                     accent_color=p["accent"], base_color=ANCHOR_DARK, size=16)
# AI DAY logo (real Figma export — multi-color treatment)
add_image(s, LOGO_AIDAY["magenta"], Inches(0.4), Inches(1.3),
          width=Inches(8.0), height=Inches(2.4))
# Ring-badge
add_ring_badge(s, int(SW * 0.78), int(SH * 0.62),
               rings=[ANCHOR_DARK, ANCHOR_LIGHT, ANCHOR_DARK, ANCHOR_LIGHT, ANCHOR_DARK],
               max_diam=Inches(2.6), stroke_w_pt=10)
# Real squiggle (yellow scribble) — bottom-left
add_image(s, SLOT_SQUIGGLE, Inches(0.0), Inches(3.6),
          width=Inches(6), height=Inches(2.4))
# AIR lockup (real Figma export)
add_image(s, LOGO_AIR["magenta"], Inches(8.2), Inches(4.2),
          width=Inches(1.5), height=Inches(1.15))
# Session title bottom-left (on top of squiggle)
add_text(s, Inches(0.4), Inches(4.8), Inches(5.5), Inches(0.35),
         "Session title goes here", font=DISPLAY_FONT,
         size=20, color=ANCHOR_LIGHT, bold=True)
add_text(s, Inches(0.4), Inches(5.15), Inches(5.5), Inches(0.25),
         "HOST PRESENTER NAME", font=MONO_FONT,
         size=10, color=ANCHOR_LIGHT, bold=False)

# ============================================================
# SLIDE 2 — Section divider (Violet) — top-edge curve
# ============================================================
p = PALETTES["violet"]
s = prs.slides.add_slide(BLANK)
set_bg(s, p["dominant"])
# Real top-edge curve (recolored to accent yellow — but png is yellow already so it works)
add_image(s, SLOT_TOP_CURVE, Inches(-1), Inches(-0.5),
          width=Inches(12), height=Inches(3))
add_eyebrow(s, Inches(0.4), Inches(0.4), "01 — CHAPTER", color=p["accent"])
add_text_runs(s, Inches(0.4), Inches(1.6), Inches(9.2), Inches(2.4),
              runs=[("Why ", ANCHOR_LIGHT), ("now", p["accent"]), (".", ANCHOR_LIGHT)],
              font=DISPLAY_FONT, size=180, bold=True)
add_text(s, Inches(0.4), Inches(4.95), Inches(8), Inches(0.3),
         "THE STATE OF AI AT THE WEATHER COMPANY",
         font=MONO_FONT, size=11, color=ANCHOR_LIGHT, bold=False)

# ============================================================
# SLIDE 3 — Text-heavy (Light field)
# ============================================================
p = PALETTES["magenta"]
s = prs.slides.add_slide(BLANK)
set_bg(s, ANCHOR_LIGHT)
add_eyebrow(s, Inches(0.4), Inches(0.4), "01 — STRATEGY", color=p["dominant"])
add_text(s, Inches(0.4), Inches(0.85), Inches(9.2), Inches(0.85),
         "A three-year AI commitment", font=DISPLAY_FONT,
         size=44, color=ANCHOR_DARK, bold=True)
add_rect(s, Inches(0.4), Inches(1.85), Inches(0.6), Inches(0.04), p["accent"])
body = (
    "By 2029, every team will have AI in their daily workflow — not as a "
    "novelty, but as core infrastructure. The investment spans three tracks: "
    "tooling, training, and culture.\n\n"
    "This is about reaching further than we ever could before — with the "
    "same people, working faster, with less friction."
)
add_text(s, Inches(0.4), Inches(2.2), Inches(5.4), Inches(3),
         body, font=DISPLAY_FONT, size=14, color=ANCHOR_DARK)
STAT_FILL = PALETTES["magenta"]["dominant"]
add_rect(s, Inches(6.4), Inches(2.2), Inches(3.2), Inches(2.6), STAT_FILL)
add_text(s, Inches(6.6), Inches(2.35), Inches(2.8), Inches(0.25),
         "BY 2029", font=MONO_FONT, size=9, color=ANCHOR_DARK, bold=False)
add_text(s, Inches(6.6), Inches(2.65), Inches(2.8), Inches(1.5),
         "$40M", font=DISPLAY_FONT, size=72, color=ANCHOR_DARK, bold=True)
add_text(s, Inches(6.6), Inches(4.0), Inches(2.8), Inches(0.8),
         "invested across tooling, training, and culture.",
         font=DISPLAY_FONT, size=12, color=ANCHOR_DARK)

# ============================================================
# SLIDE 4 — Image + text (Cyan) with converging pattern
# ============================================================
p = PALETTES["cyan"]
s = prs.slides.add_slide(BLANK)
set_bg(s, ANCHOR_LIGHT)
add_rect(s, Inches(0), Inches(0), Inches(5), SH, p["dominant"])
# Converging pattern overlay (anchor-dark on cyan)
add_image(s, SLOT_CONVERGING_DARK, Inches(-0.5), Inches(0),
          width=Inches(6), height=Inches(3))
add_rect(s, Inches(0.6), Inches(1.0), Inches(3.8), Inches(3.6), ANCHOR_LIGHT)
add_text(s, Inches(0.6), Inches(2.6), Inches(3.8), Inches(0.3),
         "[ image / demo ]", font=MONO_FONT, size=10,
         color=ANCHOR_DARK, align=PP_ALIGN.CENTER)
add_eyebrow(s, Inches(5.4), Inches(0.5), "02 — DEMO", color=p["dominant"])
add_text_runs(s, Inches(5.4), Inches(0.95), Inches(4.4), Inches(1.6),
              runs=[("Forecasting,\nrewritten ", ANCHOR_DARK),
                    ("by AI", p["dominant"]), (".", ANCHOR_DARK)],
              font=DISPLAY_FONT, size=38, bold=True)
add_rect(s, Inches(5.4), Inches(3.0), Inches(0.5), Inches(0.04), p["accent"])
add_text(s, Inches(5.4), Inches(3.2), Inches(4.4), Inches(2),
         "A walkthrough of how the forecasting team is using AI to compress "
         "two-day workflows into 15 minutes — without losing meteorological "
         "accuracy.",
         font=DISPLAY_FONT, size=12, color=ANCHOR_DARK)
add_text(s, Inches(5.4), Inches(5.0), Inches(4), Inches(0.25),
         "HOST PRESENTER NAME", font=MONO_FONT, size=9, color=ANCHOR_DARK, bold=False)

# ============================================================
# SLIDE 5 — Two-column comparison
# ============================================================
p = PALETTES["magenta"]
s = prs.slides.add_slide(BLANK)
set_bg(s, ANCHOR_LIGHT)
add_eyebrow(s, Inches(0.4), Inches(0.4), "03 — COMPARISON", color=ANCHOR_DARK)
add_text(s, Inches(0.4), Inches(0.85), Inches(9), Inches(0.85),
         "Before / After AI", font=DISPLAY_FONT,
         size=42, color=ANCHOR_DARK, bold=True)
add_rect(s, Inches(0.4), Inches(1.85), Inches(0.6), Inches(0.04), p["accent"])
add_rect(s, Inches(0.4), Inches(2.2), Inches(4.5), Inches(2.9), DARK_BOX_FILL)
add_text(s, Inches(0.6), Inches(2.35), Inches(4), Inches(0.25),
         "BEFORE", font=MONO_FONT, size=10, color=p["accent"], bold=False)
add_text(s, Inches(0.6), Inches(2.65), Inches(4), Inches(1.1),
         "2 days", font=DISPLAY_FONT, size=60, color=ANCHOR_LIGHT, bold=True)
add_text(s, Inches(0.6), Inches(3.85), Inches(4), Inches(1.2),
         "Manual data review across six systems, with multiple handoffs and "
         "reconciliation steps.",
         font=DISPLAY_FONT, size=12, color=ANCHOR_LIGHT)
add_rect(s, Inches(5.1), Inches(2.2), Inches(4.5), Inches(2.9), p["dominant"])
add_text(s, Inches(5.3), Inches(2.35), Inches(4), Inches(0.25),
         "AFTER", font=MONO_FONT, size=10, color=ANCHOR_DARK, bold=False)
add_text(s, Inches(5.3), Inches(2.65), Inches(4), Inches(1.1),
         "15 min", font=DISPLAY_FONT, size=60, color=ANCHOR_DARK, bold=True)
add_text(s, Inches(5.3), Inches(3.85), Inches(4), Inches(1.2),
         "Unified agent surface, single review, automated reconciliation. "
         "Same accuracy.",
         font=DISPLAY_FONT, size=12, color=ANCHOR_DARK)

# ============================================================
# SLIDE 6 — Two-column text (light, violet eyebrow)
# ============================================================
p = PALETTES["violet"]
s = prs.slides.add_slide(BLANK)
set_bg(s, ANCHOR_LIGHT)
add_eyebrow(s, Inches(0.4), Inches(0.4), "04 — DETAIL", color=p["dominant"])
add_text(s, Inches(0.4), Inches(0.85), Inches(9), Inches(0.85),
         "Two perspectives, one direction", font=DISPLAY_FONT,
         size=36, color=ANCHOR_DARK, bold=True)
add_rect(s, Inches(0.4), Inches(1.75), Inches(0.6), Inches(0.04), p["accent"])

col_top = Inches(2.15)
col_w = Inches(4.4)
col_h = Inches(2.9)
col1_x = Inches(0.4)
col2_x = Inches(5.2)

add_text(s, col1_x, col_top, col_w, Inches(0.3),
         "WHAT WE'RE BUILDING", font=MONO_FONT, size=10,
         color=p["dominant"], bold=False)
add_text(s, col1_x, col_top + Inches(0.4), col_w, Inches(0.5),
         "Tooling first.", font=DISPLAY_FONT, size=22,
         color=ANCHOR_DARK, bold=True)
add_text(s, col1_x, col_top + Inches(1.0), col_w, col_h - Inches(1.0),
         "Lorem ipsum dolor sit amet, consectetur adipiscing elit. Sed do "
         "eiusmod tempor incididunt ut labore et dolore magna aliqua. Ut "
         "enim ad minim veniam, quis nostrud exercitation ullamco laboris "
         "nisi ut aliquip ex ea commodo consequat.",
         font=DISPLAY_FONT, size=12, color=ANCHOR_DARK)

add_text(s, col2_x, col_top, col_w, Inches(0.3),
         "HOW WE'LL GET THERE", font=MONO_FONT, size=10,
         color=p["dominant"], bold=False)
add_text(s, col2_x, col_top + Inches(0.4), col_w, Inches(0.5),
         "Training second.", font=DISPLAY_FONT, size=22,
         color=ANCHOR_DARK, bold=True)
add_text(s, col2_x, col_top + Inches(1.0), col_w, col_h - Inches(1.0),
         "Duis aute irure dolor in reprehenderit in voluptate velit esse "
         "cillum dolore eu fugiat nulla pariatur. Excepteur sint occaecat "
         "cupidatat non proident, sunt in culpa qui officia deserunt mollit "
         "anim id est laborum.",
         font=DISPLAY_FONT, size=12, color=ANCHOR_DARK)

# ============================================================
# SLIDE 7 — Quote (Magenta) with real squiggle
# ============================================================
p = PALETTES["magenta"]
s = prs.slides.add_slide(BLANK)
set_bg(s, p["dominant"])
# Real squiggle as decoration
add_image(s, SLOT_SQUIGGLE, Inches(4), Inches(-0.5),
          width=Inches(6.5), height=Inches(2.6))
add_text(s, Inches(0.4), Inches(0.5), Inches(1.2), Inches(1.5),
         '"', font=DISPLAY_FONT, size=160, color=p["accent"], bold=True)
add_text_runs(s, Inches(0.4), Inches(2.0), Inches(9.2), Inches(2.7),
              runs=[("The bar isn't whether AI ", ANCHOR_LIGHT),
                    ("can", p["accent"]),
                    (" do this. It's whether we can move fast enough to use it "
                     "before it becomes obvious.", ANCHOR_LIGHT)],
              font=DISPLAY_FONT, size=32, bold=True)
add_rect(s, Inches(0.4), Inches(4.85), Inches(0.5), Inches(0.04), p["accent"])
add_text(s, Inches(0.4), Inches(4.95), Inches(8), Inches(0.3),
         "EXECUTIVE LEADERSHIP — THE WEATHER COMPANY",
         font=MONO_FONT, size=10, color=ANCHOR_LIGHT, bold=False)

# ============================================================
# SLIDE 8 — Blank + frame
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, ANCHOR_LIGHT)
add_eyebrow(s, Inches(0.4), Inches(0.4), "LIVE DEMO", color=ANCHOR_DARK)
acc = PALETTES["magenta"]["dominant"]
bracket = Inches(0.18)
add_rect(s, Inches(0.4), Inches(0.85), bracket, bracket, acc, shadow=False)
add_rect(s, SW - Inches(0.4) - bracket, Inches(0.85), bracket, bracket, acc, shadow=False)
add_rect(s, Inches(0.4), SH - Inches(0.6) - bracket, bracket, bracket, acc, shadow=False)
add_rect(s, SW - Inches(0.4) - bracket, SH - Inches(0.6) - bracket,
         bracket, bracket, acc, shadow=False)
add_text(s, Inches(0.4), Inches(2.7), Inches(9.2), Inches(0.3),
         "[ demo area — clear, no chrome ]",
         font=MONO_FONT, size=10, color=ANCHOR_DARK, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 9 — Big stat (Cyan)
# ============================================================
p = PALETTES["cyan"]
s = prs.slides.add_slide(BLANK)
set_bg(s, p["dominant"])
add_eyebrow(s, Inches(0.4), Inches(0.4), "04 — MOMENTUM", color=p["accent"])
add_text_runs(s, Inches(0.4), Inches(0.95), Inches(9.2), Inches(3.4),
              runs=[("9", ANCHOR_LIGHT), ("8", p["accent"]), ("%", ANCHOR_LIGHT)],
              font=DISPLAY_FONT, size=260, bold=True)
add_rect(s, Inches(0.4), Inches(4.5), Inches(0.6), Inches(0.04), p["accent"])
add_text(s, Inches(0.4), Inches(4.65), Inches(9.2), Inches(0.4),
         "Weather Company employees converging on AI Day",
         font=DISPLAY_FONT, size=20, color=ANCHOR_LIGHT, bold=True)
add_text(s, Inches(0.4), Inches(5.15), Inches(9.2), Inches(0.3),
         "FIVE GLOBAL LOCATIONS — ONE COMMITMENT",
         font=MONO_FONT, size=10, color=ANCHOR_LIGHT, bold=False)

# ============================================================
# SLIDE 10 — Three-up cards (palette showcase)
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, ANCHOR_LIGHT)
add_eyebrow(s, Inches(0.4), Inches(0.4), "05 — TRACKS", color=ANCHOR_DARK)
add_text(s, Inches(0.4), Inches(0.85), Inches(9), Inches(0.85),
         "Three tracks. One mandate.", font=DISPLAY_FONT,
         size=36, color=ANCHOR_DARK, bold=True)
add_rect(s, Inches(0.4), Inches(1.75), Inches(0.6), Inches(0.04),
         PALETTES["magenta"]["accent"], shadow=False)
tracks = [
    ("01", "TOOLING",
     "Claude, Copilot, and internal agents to every team by Q3.", "magenta"),
    ("02", "TRAINING",
     "Hands-on workshops, prompt clinics, ambassador-led sessions.", "cyan"),
    ("03", "CULTURE",
     "Make AI use visible, rewarded, and the default expectation.", "violet"),
]
col_w = Inches(3.06)
col_h = Inches(3.3)
gap = Inches(0.1)
col_top = Inches(2.05)
for i, (num, label, body, palette_key) in enumerate(tracks):
    pp = PALETTES[palette_key]
    x = Inches(0.4) + i * (col_w + gap)
    add_rect(s, x, col_top, col_w, col_h, pp["dominant"], shadow=False)
    add_text(s, x + Inches(0.25), col_top + Inches(0.2), Inches(2.5), Inches(1),
             num, font=DISPLAY_FONT, size=66, color=pp["accent"], bold=True)
    add_text(s, x + Inches(0.25), col_top + Inches(1.5), Inches(2.6), Inches(0.3),
             label, font=MONO_FONT, size=10, color=ANCHOR_LIGHT, bold=False)
    add_text(s, x + Inches(0.25), col_top + Inches(1.85), Inches(2.6), Inches(1.3),
             body, font=DISPLAY_FONT, size=12, color=ANCHOR_LIGHT)

# ============================================================
# SLIDE 11 — Hero (Cyan/royal blue) with real squiggle
# ============================================================
p = PALETTES["cyan"]
s = prs.slides.add_slide(BLANK)
set_bg(s, p["dominant"])
add_rect(s, Inches(0.4), Inches(0.4), Inches(5.5), Inches(4.85), DARK_BOX_FILL)
# Piano-stripes pattern (light variant) inside the dark box — top portion
add_image(s, SLOT_STRIPES_LIGHT, Inches(0.4), Inches(0.4),
          width=Inches(5.5), height=Inches(1.3))
add_ring_badge(s, int(Inches(3.15)), int(Inches(2.85)),
               rings=[ANCHOR_LIGHT, DARK_BOX_FILL, ANCHOR_LIGHT, DARK_BOX_FILL, ANCHOR_LIGHT],
               max_diam=Inches(3.6), stroke_w_pt=8)
add_rect(s, Inches(2.05), Inches(1.65), Inches(2.2), Inches(2.2),
         ANCHOR_LIGHT, line_color=ANCHOR_LIGHT)
add_text(s, Inches(2.05), Inches(2.6), Inches(2.2), Inches(0.3),
         "[ portrait ]", font=MONO_FONT, size=10,
         color=ANCHOR_DARK, align=PP_ALIGN.CENTER)
# Squiggle decoration
add_image(s, SLOT_SQUIGGLE, Inches(5.5), Inches(3.0),
          width=Inches(4.5), height=Inches(1.8))
add_eyebrow(s, Inches(6.2), Inches(0.5), "06 — FEATURED", color=ANCHOR_DARK)
add_text_runs(s, Inches(6.2), Inches(0.95), Inches(3.6), Inches(2.2),
              runs=[("Ambassador\n", ANCHOR_DARK),
                    ("spotlight", p["accent"]), (".", ANCHOR_DARK)],
              font=DISPLAY_FONT, size=42, bold=True)

# ============================================================
# SLIDE 12 — Closing — real quarter-circles pattern + QR + contact
# ============================================================
p = PALETTES["magenta"]
s = prs.slides.add_slide(BLANK)
set_bg(s, p["dominant"])
# Quarter-circles pattern bottom
add_image(s, SLOT_QUARTER_CIRCLES, Inches(-1), Inches(2.8),
          width=Inches(12), height=Inches(5))
# Squiggle accent — upper-right
add_image(s, SLOT_SQUIGGLE, Inches(4.5), Inches(0.3),
          width=Inches(5.5), height=Inches(2.3))
# Date
add_date_with_accent(s, Inches(0.4), Inches(0.4),
                     accent_color=p["accent"], base_color=ANCHOR_DARK, size=16)
# Big closing message — moved up to make room for next-steps + QR
add_text_runs(s, Inches(0.4), Inches(0.95), Inches(9.2), Inches(1.6),
              runs=[("See you\n", ANCHOR_DARK),
                    ("June 9", p["accent"]), (".", ANCHOR_DARK)],
              font=DISPLAY_FONT, size=110, bold=True)
# Next-steps line
add_rect(s, Inches(0.4), Inches(3.55), Inches(0.5), Inches(0.04),
         p["accent"], shadow=False)
add_text(s, Inches(0.4), Inches(3.7), Inches(5.5), Inches(0.3),
         "NEXT STEPS", font=MONO_FONT, size=10,
         color=ANCHOR_DARK, bold=False)
add_text(s, Inches(0.4), Inches(4.0), Inches(5.5), Inches(0.9),
         "RSVP via the QR. Add the date to your calendar. "
         "Bring one question worth asking out loud.",
         font=DISPLAY_FONT, size=14, color=ANCHOR_DARK)
# Contact line
add_text(s, Inches(0.4), Inches(5.0), Inches(5.5), Inches(0.25),
         "QUESTIONS  ⟶  #PROJECT-AIR ON SLACK", font=MONO_FONT,
         size=10, color=ANCHOR_DARK, bold=False)
# QR code (real if file exists; else placeholder)
qr_x, qr_y, qr_size = Inches(7.7), Inches(3.55), Inches(1.5)
if os.path.exists(QR_RSVP):
    add_image(s, QR_RSVP, qr_x, qr_y, width=qr_size, height=qr_size)
else:
    add_rect(s, qr_x, qr_y, qr_size, qr_size, ANCHOR_LIGHT,
             line_color=ANCHOR_DARK, shadow=False)
    add_text(s, qr_x, qr_y + Inches(0.6), qr_size, Inches(0.3),
             "[ QR ]", font=MONO_FONT, size=11,
             color=ANCHOR_DARK, bold=True, align=PP_ALIGN.CENTER)
add_text(s, qr_x, qr_y + qr_size + Inches(0.05), qr_size, Inches(0.2),
         "SCAN TO RSVP", font=MONO_FONT, size=8,
         color=ANCHOR_DARK, bold=False, align=PP_ALIGN.CENTER)
# AIR lockup
add_image(s, LOGO_AIR["magenta"], Inches(8.2), Inches(0.4),
          width=Inches(1.4), height=Inches(1.1))

# ============================================================
# SLIDE 13 — How to use
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, ANCHOR_LIGHT)
add_eyebrow(s, Inches(0.4), Inches(0.4), "GUIDE", color=ANCHOR_DARK)
add_text(s, Inches(0.4), Inches(0.85), Inches(9), Inches(0.85),
         "How to use this template", font=DISPLAY_FONT,
         size=32, color=ANCHOR_DARK, bold=True)
add_rect(s, Inches(0.4), Inches(1.7), Inches(0.6), Inches(0.04),
         PALETTES["magenta"]["accent"], shadow=False)
add_text(s, Inches(0.4), Inches(2.0), Inches(5), Inches(0.25),
         "3 PALETTES — SWAP DOMINANT + ACCENT PER SLIDE",
         font=MONO_FONT, size=8, color=ANCHOR_DARK, bold=False)
palette_keys = ["magenta", "cyan", "violet"]
sw = Inches(0.45)
sw_top = Inches(2.3)
for i, key in enumerate(palette_keys):
    pp = PALETTES[key]
    x = Inches(0.4) + i * (sw + Inches(0.55))
    add_rect(s, x, sw_top, sw, sw, pp["dominant"], shadow=False)
    add_rect(s, x + Inches(0.25), sw_top + Inches(0.25), Inches(0.2), Inches(0.2),
             pp["accent"], shadow=False)
    add_text(s, x - Inches(0.05), sw_top + sw + Inches(0.05), Inches(0.6), Inches(0.18),
             key.upper(), font=MONO_FONT, size=7, color=ANCHOR_DARK, bold=False)
add_text(s, Inches(0.4), Inches(3.5), Inches(5), Inches(0.25),
         "TYPOGRAPHY", font=MONO_FONT, size=8, color=ANCHOR_DARK, bold=False)
add_text(s, Inches(0.4), Inches(3.8), Inches(5), Inches(0.5),
         "Aa", font=DISPLAY_FONT, size=48, color=ANCHOR_DARK, bold=True)
add_text(s, Inches(1.3), Inches(3.9), Inches(4), Inches(0.3),
         "Fira Sans Extra Condensed", font=DISPLAY_FONT,
         size=12, color=ANCHOR_DARK, bold=True)
add_text(s, Inches(1.3), Inches(4.15), Inches(4), Inches(0.3),
         "Display + body — all headings, body copy",
         font=DISPLAY_FONT, size=10, color=ANCHOR_DARK)
add_text(s, Inches(0.4), Inches(4.65), Inches(5), Inches(0.5),
         "Aa", font=MONO_FONT, size=32, color=ANCHOR_DARK, bold=True)
add_text(s, Inches(1.3), Inches(4.7), Inches(4), Inches(0.3),
         "Fira Code", font=DISPLAY_FONT,
         size=12, color=ANCHOR_DARK, bold=True)
add_text(s, Inches(1.3), Inches(4.95), Inches(4), Inches(0.3),
         "Mono — eyebrows, captions, labels",
         font=DISPLAY_FONT, size=10, color=ANCHOR_DARK)
add_text(s, Inches(6), Inches(2.0), Inches(4), Inches(0.25),
         "RULES", font=MONO_FONT, size=8, color=ANCHOR_DARK, bold=False)
rules = [
    "1.  One palette per slide (or section).",
    "2.  Anchor-dark + anchor-light are constants.",
    "3.  Accent-primary highlights ONE word/digit, never all.",
    "4.  No team-member names anywhere in the deck.",
    "5.  Speaker notes live in the notes pane, not on slide.",
    "6.  Lockup always bottom-right.",
]
for i, line in enumerate(rules):
    add_text(s, Inches(6), Inches(2.35) + i * Inches(0.42),
             Inches(4), Inches(0.4),
             line, font=DISPLAY_FONT, size=11, color=ANCHOR_DARK)

# ============================================================
# SLIDE 14 — Build your deck (Workflow guide)
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, ANCHOR_LIGHT)
add_eyebrow(s, Inches(0.4), Inches(0.4), "GUIDE — WORKFLOW", color=ANCHOR_DARK)
add_text(s, Inches(0.4), Inches(0.85), Inches(9), Inches(0.85),
         "Build your deck in 4 steps", font=DISPLAY_FONT,
         size=32, color=ANCHOR_DARK, bold=True)
add_rect(s, Inches(0.4), Inches(1.7), Inches(0.6), Inches(0.04),
         PALETTES["magenta"]["accent"], shadow=False)

steps = [
    ("01", "Duplicate a slide",
     "Right-click any slide in the left pane → Duplicate Slide. "
     "Or press Cmd+D (Mac) / Ctrl+D (Windows).",
     "magenta"),
    ("02", "Replace the text",
     "Click any text to edit. Keep the rhythm: eyebrow → headline → "
     "body → caption. Don't change fonts or colors — they're locked.",
     "cyan"),
    ("03", "Swap in your image",
     "Right-click any [placeholder] block → Change Picture. The frame "
     "shape stays — your image fills it. Use high-res sources.",
     "violet"),
    ("04", "Save your own file",
     "File → Save As → name it for your session. Never edit this master "
     "template directly. Add your speaker notes in the Notes pane.",
     "magenta"),
]
# Two rows of two steps (2x2 grid)
step_w = Inches(4.55)
step_h = Inches(1.30)
gap_x = Inches(0.2)
gap_y = Inches(0.15)
start_top = Inches(2.05)
for i, (num, label, body, palette_key) in enumerate(steps):
    pp = PALETTES[palette_key]
    col = i % 2
    row = i // 2
    x = Inches(0.4) + col * (step_w + gap_x)
    y = start_top + row * (step_h + gap_y)
    # number block
    add_rect(s, x, y, Inches(1), step_h, pp["dominant"], shadow=False)
    add_text(s, x, y + Inches(0.22), Inches(1), Inches(0.9),
             num, font=DISPLAY_FONT, size=44, color=pp["accent"],
             bold=True, align=PP_ALIGN.CENTER)
    # text block
    add_rect(s, x + Inches(1), y, step_w - Inches(1), step_h, ANCHOR_LIGHT,
             line_color=ANCHOR_DARK, shadow=False)
    add_text(s, x + Inches(1.15), y + Inches(0.12), step_w - Inches(1.2),
             Inches(0.3), label,
             font=DISPLAY_FONT, size=15, color=ANCHOR_DARK, bold=True)
    add_text(s, x + Inches(1.15), y + Inches(0.45), step_w - Inches(1.2),
             Inches(0.8), body,
             font=DISPLAY_FONT, size=10, color=ANCHOR_DARK)

# Footer tip
add_text(s, Inches(0.4), Inches(5.25), Inches(9.2), Inches(0.3),
         "STUCK?  PING #PROJECT-AIR-CREATIVE ON SLACK — OR PICK THE "
         "CLOSEST LAYOUT AND ADAPT.",
         font=MONO_FONT, size=9, color=ANCHOR_DARK, bold=False,
         align=PP_ALIGN.CENTER)

# --- speaker notes (template placeholders) ---
SLIDE_NOTES = [
    # 1 — Title
    "TITLE SLIDE\n\n"
    "Open with a hook (not a self-intro). 30-45 seconds max.\n"
    "Suggested beats:\n"
    "• Name the moment — why this room, why today\n"
    "• Land the session title in one breath\n"
    "• Tee up the date as a stake in the ground (June 9, 2026)\n\n"
    "EDIT: Replace 'Session title goes here' and 'HOST PRESENTER NAME'.",

    # 2 — Section divider
    "SECTION DIVIDER — 'Why now'\n\n"
    "Use to mark a chapter break. Talk for 60-90 seconds.\n"
    "Frame the urgency: what changed, what's at stake, what we're "
    "responding to. Don't list facts here — set the question that "
    "the next 2-3 slides will answer.\n\n"
    "EDIT: Change eyebrow ('01 — CHAPTER') and headline to match your section.",

    # 3 — Text-heavy
    "STRATEGY / COMMITMENT SLIDE\n\n"
    "Use when the audience needs a written reference point.\n"
    "Read the headline aloud, then paraphrase the body — don't read it "
    "verbatim. Let the stat block on the right do the heavy lifting.\n\n"
    "EDIT: Headline, body copy, stat block ($40M / BY 2029 / caption). "
    "Keep stat box to one short number + one short qualifier.",

    # 4 — Image + text
    "IMAGE + TEXT (HALF/HALF)\n\n"
    "Use for any 'show + tell' moment — product shot, screenshot, photo.\n"
    "Talk to the image first, then bridge to the headline.\n\n"
    "EDIT: Replace [placeholder] image (right-click → Change Picture). "
    "Update headline runs and the HOST PRESENTER NAME line.",

    # 5 — Comparison
    "BEFORE / AFTER COMPARISON\n\n"
    "Strongest when the contrast is dramatic — minutes vs. days, "
    "manual vs. automated. Read BEFORE state, pause, then deliver AFTER "
    "as the payoff.\n\n"
    "EDIT: BEFORE box (deep violet) and AFTER box (hot pink). Swap the "
    "numbers and the one-line explanation under each.",

    # 6 — Two-column text
    "TWO-COLUMN TEXT\n\n"
    "Use for parallel ideas — pros/cons, two principles, two examples. "
    "Don't overload; ~50 words per column max.\n\n"
    "EDIT: Replace both column headlines and body copy. Keep the eyebrow "
    "tag for navigation context.",

    # 7 — Quote
    "QUOTE / CALLOUT\n\n"
    "Use a real quote with attribution, or a guiding principle.\n"
    "Pause after reading. Let it land before moving on.\n\n"
    "EDIT: Replace quote text and attribution line.",

    # 8 — Blank + frame (Live demo)
    "LIVE DEMO / FRAME SLIDE\n\n"
    "Switch to your demo here. The bracket marks ([ ]) are visual "
    "guides — drop your screen share / embedded video inside them.\n\n"
    "EDIT: Replace placeholder text. If you're not demoing live, paste "
    "a screenshot or short video into the frame area.",

    # 9 — Big stat
    "BIG STAT SLIDE\n\n"
    "One number, one claim. Say the number, pause, then explain it.\n"
    "Strongest when the stat is specific and recent.\n\n"
    "EDIT: Number (currently '98%'), the caption underneath, and source "
    "line at the bottom. Always cite your source.",

    # 10 — Three-up cards
    "THREE-UP CARDS (TRACKS / PILLARS)\n\n"
    "Use for trios — three tracks, three principles, three takeaways. "
    "Walk left → middle → right. Don't try to cover everything; pick the "
    "headline for each.\n\n"
    "EDIT: Three card labels and body copy. Numbering (01/02/03) is "
    "decorative — keep it.",

    # 11 — Hero / Spotlight
    "FEATURED HERO SLIDE\n\n"
    "Use to spotlight a person, project, or team without naming "
    "individuals on the deck itself (per Mark's directive). Speak to "
    "their work, not their title.\n\n"
    "EDIT: Replace [portrait] placeholder. Update the headline runs.",

    # 12 — Closing
    "CLOSING SLIDE\n\n"
    "Three things to leave them with:\n"
    "1) The date (June 9) — said out loud\n"
    "2) The next step (RSVP via QR / Slack)\n"
    "3) Where to follow up (#project-air on Slack)\n\n"
    "Don't read the QR — point to it. Hold the slide for ~10 seconds so "
    "people can scan.\n\n"
    "EDIT: Contact line, next-steps copy. Replace QR placeholder by saving "
    "qr-rsvp.png to assets/templates/slot-elements/png/.",

    # 13 — How to use
    "TEMPLATE GUIDE — not part of your presentation.\n\n"
    "This slide explains the system to anyone editing the deck. "
    "Leave it in the master file; delete it from your own copies.\n\n"
    "Palette swatches show the three pairings (magenta / cyan / violet). "
    "Type sample shows the two fonts (Fira Sans Extra Condensed for "
    "display, Fira Code for mono labels).",

    # 14 — Build your deck
    "TEMPLATE GUIDE — workflow steps.\n\n"
    "Four-step workflow for new presenters:\n"
    "01 Duplicate a slide  →  02 Replace text  →  03 Swap image  →  "
    "04 Save As\n\n"
    "Delete this slide from your own copy. Pair with slide 13 (palette + "
    "type guide) for the full how-to.",
]

for i, note in enumerate(SLIDE_NOTES):
    add_notes(prs.slides[i], note)


# --- save ---
out = "/Users/lee.payne@weather.com/Dev/project-air-creative/assets/templates/Project_AIR_Template_v13.pptx"
prs.save(out)


# --- patch theme so palette appears in PowerPoint color picker ---
import zipfile, shutil, re, os, tempfile

THEME_COLORS = {
    "dk2":      "292929",
    "lt2":      "F0F0EB",
    "accent1":  "FB00FF",  # magenta dominant
    "accent2":  "F4DC52",  # magenta/violet accent (warm yellow)
    "accent3":  "0062FF",  # cyan dominant (royal blue)
    "accent4":  "67FAE0",  # cyan accent (mint/aqua)
    "accent5":  "46125B",  # violet dominant
    "accent6":  "F4DC52",  # violet accent
    "hlink":    "0062FF",
    "folHlink": "46125B",
}


def patch_theme(pptx_path, colors):
    fd, tmp = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)
    with zipfile.ZipFile(pptx_path, "r") as zin, \
         zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.startswith("ppt/theme/theme") and item.filename.endswith(".xml"):
                xml = data.decode("utf-8")
                for tag, hex_color in colors.items():
                    xml = re.sub(
                        rf'<a:{tag}>\s*<a:(?:srgbClr|sysClr)[^/]*/>\s*</a:{tag}>',
                        f'<a:{tag}><a:srgbClr val="{hex_color}"/></a:{tag}>',
                        xml,
                    )
                data = xml.encode("utf-8")
            zout.writestr(item, data)
    shutil.move(tmp, pptx_path)


patch_theme(out, THEME_COLORS)

print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
