"""
Extend R1B.pptx with 2 new palettes (Amber + Forest), saving as R1C.

- Update slide 13 "How to use" — add Amber + Forest swatch pairs, "3 PALETTES" -> "5 PALETTES"
- Add slide 15 — Amber palette demo (section divider style)
- Add slide 16 — Forest palette demo (stat style)
- Patch theme1.xml to expose the 5 dominants in PowerPoint's color picker
"""
import copy, os, re, shutil, zipfile
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree

SRC = "/Users/lee.payne@weather.com/Dev/project-air-creative/assets/templates/AIDAY-PPT-Template-R1B.pptx"
OUT = "/Users/lee.payne@weather.com/Dev/project-air-creative/assets/templates/AIDAY-PPT-Template-R1C.pptx"

# --- Tokens
ANCHOR_DARK = RGBColor(0x29, 0x29, 0x29)
ANCHOR_LIGHT = RGBColor(0xF0, 0xF0, 0xEB)

AMBER_DOM = RGBColor(0xFF, 0x95, 0x00)
AMBER_ACCENT = RGBColor(0xBC, 0x11, 0x00)
FOREST_DOM = RGBColor(0x1F, 0x7A, 0x4D)
FOREST_ACCENT = RGBColor(0x0D, 0x14, 0x2A)

FONT_DISPLAY = "Fira Sans Extra Condensed"
FONT_MONO = "Fira Code"


def no_shadow(shape):
    spPr = shape._element.spPr
    for el in spPr.findall(qn("a:effectLst")):
        spPr.remove(el)
    eff = etree.SubElement(spPr, qn("a:effectLst"))


def set_solid_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def add_rect(slide, l, t, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    set_solid_fill(s, fill)
    if line is None:
        s.line.fill.background()
    no_shadow(s)
    return s


def add_text(slide, l, t, w, h, text, *, font=FONT_DISPLAY, size=14, bold=False,
             color=ANCHOR_DARK, align=None, tracking=None):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    if tracking is not None:
        rPr = r._r.get_or_add_rPr()
        rPr.set("spc", str(tracking))
    return tb


# ============================================================
# Open R1B
# ============================================================
prs = Presentation(SRC)

# ============================================================
# SLIDE 13 — update palettes section
# ============================================================
slide13 = prs.slides[12]

# Update header "3 PALETTES" -> "5 PALETTES"
for sh in slide13.shapes:
    if sh.has_text_frame and "3 PALETTES" in sh.text_frame.text:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if "3 PALETTES" in r.text:
                    r.text = r.text.replace("3 PALETTES", "5 PALETTES")

# Existing palette swatches sit at L=0.40, 1.40, 2.40 (magenta/cyan/violet).
# Add Amber at L=3.40 and Forest at L=4.40, matching the existing dimensions:
# Big square: 0.45 x 0.45 at L+0.00, T=2.30 (dominant)
# Small accent: 0.20 x 0.20 at L+0.25, T=2.55 (accent overlap)
# Label: at L-0.05, T=2.80, W=0.60, H=0.18

NOTES_BODY_XML = """<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr>
    <p:cNvPr id="3" name="Notes Placeholder 2"/>
    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
    <p:nvPr><p:ph type="body" sz="quarter" idx="3"/></p:nvPr>
  </p:nvSpPr>
  <p:spPr/>
  <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:r><a:rPr dirty="0"/><a:t></a:t></a:r></a:p></p:txBody>
</p:sp>"""


def add_notes(slide, text):
    """Add speaker notes. If the freshly-added notes slide lacks a body placeholder, inject one."""
    ns = slide.notes_slide
    if ns.notes_text_frame is None:
        # Inject body placeholder into the notes spTree
        spTree = ns.element.find(qn("p:cSld")).find(qn("p:spTree"))
        sp = parse_xml(NOTES_BODY_XML)
        spTree.append(sp)
    ns.notes_text_frame.text = text


def add_swatch(slide, l, dom, accent, label):
    # Big dominant
    big = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(2.30), Inches(0.45), Inches(0.45))
    set_solid_fill(big, dom)
    big.line.fill.background()
    no_shadow(big)
    # Small accent overlay
    sm = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l + 0.25), Inches(2.55), Inches(0.20), Inches(0.20))
    set_solid_fill(sm, accent)
    sm.line.fill.background()
    no_shadow(sm)
    # Label
    add_text(slide, l - 0.05, 2.80, 0.60, 0.18, label,
             font=FONT_MONO, size=6, color=ANCHOR_DARK, tracking=100)

add_swatch(slide13, 3.40, AMBER_DOM, AMBER_ACCENT, "AMBER")
add_swatch(slide13, 4.40, FOREST_DOM, FOREST_ACCENT, "FOREST")

# ============================================================
# Add 2 new slides at the END (will sit after slide 14 build-deck)
# Use blank layout
# ============================================================
BLANK = prs.slide_layouts[6]


# ---------- SLIDE 15 — Amber section divider ----------
s15 = prs.slides.add_slide(BLANK)

# Full-bleed amber background
bg = s15.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.625))
set_solid_fill(bg, AMBER_DOM)
bg.line.fill.background()
no_shadow(bg)

# Eyebrow
add_text(s15, 0.40, 0.40, 4.00, 0.20, "07 — DEMO",
         font=FONT_MONO, size=9, color=ANCHOR_LIGHT, tracking=200)

# Big headline
add_text(s15, 0.40, 1.60, 9.20, 2.89, "Take it further.",
         font=FONT_DISPLAY, size=120, bold=True, color=ANCHOR_LIGHT)

# Subline
add_text(s15, 0.40, 4.65, 0.60, 0.04, "", color=AMBER_ACCENT)  # placeholder
rule = s15.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.40), Inches(4.50), Inches(0.60), Inches(0.04))
set_solid_fill(rule, AMBER_ACCENT)
rule.line.fill.background()
no_shadow(rule)

add_text(s15, 0.40, 4.65, 9.20, 0.40, "Where AI shifts from novelty to leverage",
         font=FONT_DISPLAY, size=22, color=ANCHOR_LIGHT)

add_text(s15, 0.40, 5.15, 9.20, 0.30, "PALETTE — AMBER + DEEP RED",
         font=FONT_MONO, size=9, color=ANCHOR_LIGHT, tracking=200)

# Speaker notes
add_notes(s15, (
    "TEMPLATE SLIDE — Amber palette section divider.\n\n"
    "What this is for: Use this layout to introduce a new chapter or pivot in your deck. "
    "The full-bleed amber + dark red accent is a high-energy, high-contrast moment — use sparingly.\n\n"
    "EDIT:\n"
    "- Change the eyebrow (07 — DEMO) to your section number + label.\n"
    "- Replace 'Take it further.' with your chapter headline (≤ 6 words).\n"
    "- Replace the subline with one sentence on why this section matters.\n"
    "- Keep the bottom-right lockup placement consistent with other slides."
))


# ---------- SLIDE 16 — Forest stat ----------
s16 = prs.slides.add_slide(BLANK)

# Cream background already (PPT default), so add nothing for bg

# Eyebrow
add_text(s16, 0.40, 0.40, 4.00, 0.20, "08 — IMPACT",
         font=FONT_MONO, size=9, color=ANCHOR_DARK, tracking=200)

# Big stat in Forest dominant
add_text(s16, 0.40, 0.53, 9.20, 4.38, "5",
         font=FONT_DISPLAY, size=520, bold=True, color=FOREST_DOM)

# Accent rule (forest's contrast navy)
rule = s16.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.40), Inches(4.50), Inches(0.60), Inches(0.04))
set_solid_fill(rule, FOREST_ACCENT)
rule.line.fill.background()
no_shadow(rule)

add_text(s16, 0.40, 4.65, 9.20, 0.40,
         "Offices converging on a single day, a single mandate",
         font=FONT_DISPLAY, size=22, color=ANCHOR_DARK)

add_text(s16, 0.40, 5.15, 9.20, 0.30,
         "BROOKHAVEN  /  ANDOVER  /  BIRMINGHAM  /  SEOUL  /  NEW YORK",
         font=FONT_MONO, size=9, color=FOREST_DOM, tracking=200)

# Speaker notes
add_notes(s16, (
    "TEMPLATE SLIDE — Forest palette stat slide.\n\n"
    "What this is for: Anchor a single number with weight. One stat. One subline. Nothing else competes.\n\n"
    "EDIT:\n"
    "- Replace '5' with your stat. Keep it short — one or two digits, or a tight word like 'EVERY'.\n"
    "- Update the subline below the rule with a 6–10 word descriptor.\n"
    "- The mono caption at the bottom is optional context (locations, sources, etc.).\n"
    "- Forest is the most muted of the 5 palettes — best for serious or factual moments."
))


# ============================================================
# Save .pptx
# ============================================================
prs.save(OUT)
print(f"Saved {OUT}")

# ============================================================
# Patch theme1.xml — expose 5 dominants + shared yellow accent in PowerPoint UI
# accent1: magenta, accent2: cyan, accent3: violet, accent4: amber, accent5: forest, accent6: yellow
# dk2: anchor dark, lt2: anchor light
# hlink: cyan, folHlink: violet
# ============================================================
THEME_COLORS = {
    "dk2":      "292929",
    "lt2":      "F0F0EB",
    "accent1":  "FB00FF",
    "accent2":  "0062FF",
    "accent3":  "46125B",
    "accent4":  "FF9500",
    "accent5":  "1F7A4D",
    "accent6":  "F4DC52",
    "hlink":    "0062FF",
    "folHlink": "46125B",
}

with zipfile.ZipFile(OUT) as z:
    theme_xml = z.read("ppt/theme/theme1.xml").decode("utf-8")

for slot, hex_val in THEME_COLORS.items():
    # Match <a:slot>...<a:srgbClr val="XXXXXX"/>...</a:slot> OR <a:sysClr ... />
    # Replace inner color with srgbClr
    pattern_srgb = rf'(<a:{slot}>)\s*<a:srgbClr val="[0-9A-Fa-f]{{6}}"\s*/>\s*(</a:{slot}>)'
    pattern_sys = rf'(<a:{slot}>)\s*<a:sysClr [^/]*?/>\s*(</a:{slot}>)'
    replacement = rf'\1<a:srgbClr val="{hex_val}"/>\2'
    new_xml = re.sub(pattern_srgb, replacement, theme_xml)
    if new_xml == theme_xml:
        new_xml = re.sub(pattern_sys, replacement, theme_xml)
    if new_xml == theme_xml:
        print(f"  WARN: did not match {slot}")
    theme_xml = new_xml

# Write theme back into the .pptx (zip update)
tmp = OUT + ".tmp"
with zipfile.ZipFile(OUT, "r") as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
    for item in zin.infolist():
        data = zin.read(item.filename)
        if item.filename == "ppt/theme/theme1.xml":
            data = theme_xml.encode("utf-8")
        zout.writestr(item, data)
os.replace(tmp, OUT)
print(f"Theme patched in {OUT}")
